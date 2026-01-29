#!/usr/bin/env python3
"""Generate PowerPoint presentation for MCP Context Forge."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# IBM Colors
IBM_BLUE = RGBColor(0, 98, 155)  # #00629B
IBM_DARK_BLUE = RGBColor(0, 45, 75)  # #002D4B
IBM_GRAY = RGBColor(50, 50, 50)
IBM_LIGHT_GRAY = RGBColor(244, 244, 244)
WHITE = RGBColor(255, 255, 255)
ACCENT_BLUE = RGBColor(0, 129, 198)  # #0081C6


def set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_shape(slide, text, top, font_size=44, bold=True, color=IBM_DARK_BLUE):
    """Add a title text box."""
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(1)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return shape


def add_body_text(slide, text, top, left=0.5, font_size=18, color=IBM_GRAY, width=9):
    """Add body text."""
    shape = slide.shapes.add_textbox(Inches(left), top, Inches(width), Inches(4))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return shape


def add_bullet_points(slide, bullets, top, left=0.5, font_size=18, color=IBM_GRAY):
    """Add bullet points."""
    shape = slide.shapes.add_textbox(Inches(left), top, Inches(8.5), Inches(4.5))
    tf = shape.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    return shape


def add_two_column_bullets(slide, left_title, left_bullets, right_title, right_bullets, top):
    """Add two columns of bullet points."""
    # Left column title
    left_shape = slide.shapes.add_textbox(Inches(0.5), top, Inches(4.2), Inches(0.5))
    tf = left_shape.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE

    # Left column bullets
    left_bullets_shape = slide.shapes.add_textbox(Inches(0.5), top + Inches(0.5), Inches(4.2), Inches(3))
    tf = left_bullets_shape.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = IBM_GRAY
        p.space_before = Pt(4)

    # Right column title
    right_shape = slide.shapes.add_textbox(Inches(5.2), top, Inches(4.2), Inches(0.5))
    tf = right_shape.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE

    # Right column bullets
    right_bullets_shape = slide.shapes.add_textbox(Inches(5.2), top + Inches(0.5), Inches(4.2), Inches(3))
    tf = right_bullets_shape.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = IBM_GRAY
        p.space_before = Pt(4)


def add_diagram_box(slide, text, left, top, width=2, height=0.8, fill_color=IBM_BLUE, text_color=WHITE):
    """Add a diagram box with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow_down(slide, left, top, length=0.5):
    """Add a down arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(0.3), Inches(length)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = IBM_GRAY
    shape.line.color.rgb = IBM_GRAY
    return shape


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Title Slide ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, WHITE)

    # Blue header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.3))
    header.fill.solid()
    header.fill.fore_color.rgb = IBM_BLUE
    header.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ENTERPRISE AI TOOL GOVERNANCE"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = IBM_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Model Context Protocol (MCP) & MCP Context Forge"
    p.font.size = Pt(28)
    p.font.color.rgb = IBM_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.6))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "A Production-Grade Approach to Scalable, Governed, Agentic AI Infrastructure"
    p.font.size = Pt(18)
    p.font.color.rgb = IBM_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Presenter info
    presenter_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.2))
    tf = presenter_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ruslan Idelfonso Magana Vsevolodovna"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = IBM_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Data Scientist & AI Engineer"
    p2.font.size = Pt(16)
    p2.font.color.rgb = IBM_GRAY
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "IBM Client Innovation Center Italy"
    p3.font.size = Pt(16)
    p3.font.color.rgb = IBM_GRAY
    p3.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 2: Session Objectives ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "WHAT YOU WILL LEARN TODAY", Inches(0.5))

    objectives = [
        "Why AI systems require standardized tool interfaces",
        "How Model Context Protocol (MCP) addresses this need",
        "What MCP Context Forge provides as a governance layer",
        "Enterprise capabilities: federation, multi-tenancy, security, and observability",
        "Live demonstration of tool discovery and governance"
    ]
    add_bullet_points(slide, objectives, Inches(1.5), font_size=22)

    # ========== SLIDE 3: Evolution of AI ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "THE EVOLUTION OF AI SYSTEMS", Inches(0.5))

    # Timeline boxes
    add_diagram_box(slide, "2020-2022", 0.5, 1.8, 2, 0.6, IBM_GRAY)
    add_body_text(slide, "Question Answering\nModel → Text Response", Inches(2.0), left=2.7, font_size=16)

    add_diagram_box(slide, "2023-2024", 0.5, 3.2, 2, 0.6, IBM_BLUE)
    add_body_text(slide, "Tool-Augmented AI\nModel → Tools → Actions", Inches(3.6), left=2.7, font_size=16)

    add_diagram_box(slide, "2025+", 0.5, 4.6, 2, 0.6, ACCENT_BLUE)
    add_body_text(slide, "Agentic AI Ecosystems\nAgents ↔ Services ↔ Agents", Inches(5.0), left=2.7, font_size=16)

    # ========== SLIDE 4: Enterprise Challenge ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "THE ENTERPRISE CHALLENGE", Inches(0.5))

    add_body_text(slide, "Modern AI systems depend on diverse tool ecosystems:", Inches(1.4), font_size=20)

    # Grid of tool types
    tools = ["REST APIs", "Databases", "Legacy Systems", "gRPC Services", "Internal Platforms", "Cloud Services"]
    for i, tool in enumerate(tools):
        row = i // 3
        col = i % 3
        add_diagram_box(slide, tool, 0.8 + col * 3, 2.2 + row * 1.2, 2.5, 0.8, IBM_LIGHT_GRAY, IBM_DARK_BLUE)

    add_body_text(slide, "Each requires custom integration, authentication, and error handling", Inches(5.0), font_size=18, color=IBM_GRAY)

    # ========== SLIDE 5: Pain Points ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "PAIN POINTS AT ENTERPRISE SCALE", Inches(0.5))

    pain_points = [
        ("1. TOOL SPRAWL", "No central inventory • Duplicate integrations across teams"),
        ("2. INCONSISTENT SECURITY", "Each integration handles auth differently • No unified audit trail"),
        ("3. OPERATIONAL COMPLEXITY", "No visibility into usage patterns • Cannot disable tools centrally"),
        ("4. VENDOR LOCK-IN", "Tight coupling to AI providers • Expensive to switch")
    ]

    y = 1.5
    for title, desc in pain_points:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = IBM_BLUE

        desc_shape = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.4), Inches(8.5), Inches(0.5))
        tf = desc_shape.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = IBM_GRAY

        y += 1.3

    # ========== SLIDE 6: Agentic AI Complexity ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "AGENTIC AI AMPLIFIES COMPLEXITY", Inches(0.5))

    add_two_column_bullets(
        slide,
        "Traditional AI",
        ["Human in the loop", "Fixed tool selection", "Supervised execution", "Single model focus"],
        "Agentic AI",
        ["Autonomous operation", "Dynamic tool discovery", "Multi-agent coordination", "Cross-system workflows"],
        Inches(1.5)
    )

    conclusion = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.6))
    tf = conclusion.text_frame
    p = tf.paragraphs[0]
    p.text = "Agents need DYNAMIC, GOVERNED tool discovery"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 7: Introducing MCP ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "INTRODUCING MODEL CONTEXT PROTOCOL (MCP)", Inches(0.5))

    add_body_text(slide, "An open standard that defines how AI systems interact with tools, resources, and prompts", Inches(1.4), font_size=20)

    # Architecture diagram
    add_diagram_box(slide, "AI Client\n(Claude, GPT, watsonx)", 1, 2.5, 3, 1, IBM_DARK_BLUE)
    add_diagram_box(slide, "MCP Protocol\nJSON-RPC 2.0", 4.5, 2.5, 2.5, 1, ACCENT_BLUE)
    add_diagram_box(slide, "MCP Server\n(Tools, Resources, Prompts)", 7.5, 2.5, 2, 1, IBM_BLUE)

    # Arrows
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4), Inches(2.85), Inches(0.5), Inches(0.3))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = IBM_GRAY

    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7), Inches(2.85), Inches(0.5), Inches(0.3))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = IBM_GRAY

    add_body_text(slide, "Key: Standardized interface, any implementation", Inches(4.5), font_size=18, color=IBM_BLUE)

    # ========== SLIDE 8: MCP Core Concepts ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "MCP CORE CONCEPTS", Inches(0.5))

    concepts = [
        ("TOOLS", "Executable functions with defined schemas", "search_database, send_email, run_query"),
        ("RESOURCES", "URI-addressable data sources", "file://config.json, db://users/123"),
        ("PROMPTS", "Reusable templates with placeholders", "code_review_prompt, summary_template")
    ]

    y = 1.5
    for title, desc, examples in concepts:
        # Title
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(2), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = IBM_BLUE
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Description
        desc_shape = slide.shapes.add_textbox(Inches(2.7), Inches(y), Inches(6.5), Inches(0.4))
        tf = desc_shape.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = IBM_DARK_BLUE

        # Examples
        ex_shape = slide.shapes.add_textbox(Inches(2.7), Inches(y + 0.4), Inches(6.5), Inches(0.4))
        tf = ex_shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"Examples: {examples}"
        p.font.size = Pt(14)
        p.font.color.rgb = IBM_GRAY

        y += 1.5

    # ========== SLIDE 9: Why MCP Matters ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "WHY MCP MATTERS FOR ENTERPRISE", Inches(0.5))

    benefits = [
        ("DECOUPLING", "Models and tools evolve independently"),
        ("REUSABILITY", "Build once, use across all AI systems"),
        ("VENDOR FLEXIBILITY", "Switch AI providers without rewriting integrations"),
        ("FUTURE-PROOF", "Ready for agentic workflows and multi-agent systems")
    ]

    for i, (title, desc) in enumerate(benefits):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 4.7
        y = 1.8 + row * 2.2

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.3), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = IBM_LIGHT_GRAY
        box.line.color.rgb = IBM_BLUE

        title_shape = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(3.9), Inches(0.5))
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = IBM_BLUE

        desc_shape = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(3.9), Inches(0.9))
        tf = desc_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = IBM_GRAY

    # ========== SLIDE 10: Missing Piece ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "THE MISSING PIECE: CENTRAL GOVERNANCE", Inches(0.5))

    add_body_text(slide, "MCP standardizes the interface, but enterprises need answers to:", Inches(1.4), font_size=20)

    questions = [
        "How do we manage dozens of MCP servers?",
        "How do we enforce consistent security?",
        "How do we enable/disable tools centrally?",
        "How do we provide team-based access control?",
        "How do we monitor and audit tool usage?"
    ]

    for i, q in enumerate(questions):
        q_shape = slide.shapes.add_textbox(Inches(1), Inches(2.2 + i * 0.7), Inches(8), Inches(0.6))
        tf = q_shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"?  {q}"
        p.font.size = Pt(20)
        p.font.color.rgb = IBM_DARK_BLUE

    answer = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.6))
    tf = answer.text_frame
    p = tf.paragraphs[0]
    p.text = "This is where MCP Context Forge comes in"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 11: Introducing Context Forge ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "INTRODUCING MCP CONTEXT FORGE", Inches(0.5))

    add_body_text(slide, "A production-grade gateway, proxy, and registry for MCP servers and A2A agents", Inches(1.3), font_size=20, color=IBM_DARK_BLUE)

    features = [
        "Gateway & Proxy Layer",
        "Federation across multiple services",
        "Virtual server composition",
        "Multi-transport support (SSE, WebSocket, HTTP)",
        "Admin UI for real-time management",
        "Enterprise security & observability"
    ]
    add_bullet_points(slide, features, Inches(2.2), font_size=20)

    github = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf = github.text_frame
    p = tf.paragraphs[0]
    p.text = "Open source: github.com/IBM/mcp-context-forge"
    p.font.size = Pt(18)
    p.font.color.rgb = IBM_BLUE
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 12: Architecture ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "HIGH-LEVEL ARCHITECTURE", Inches(0.5))

    # Top layer - AI Clients
    add_diagram_box(slide, "AI Clients / Agents\n(Claude, watsonx, GPT, Custom)", 2, 1.5, 6, 0.9, IBM_DARK_BLUE)

    add_arrow_down(slide, 4.85, 2.45, 0.4)

    # Middle layer - Context Forge
    forge_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3), Inches(8), Inches(1.8))
    forge_box.fill.solid()
    forge_box.fill.fore_color.rgb = ACCENT_BLUE
    forge_box.line.color.rgb = IBM_BLUE

    forge_title = slide.shapes.add_textbox(Inches(1), Inches(3.1), Inches(8), Inches(0.5))
    tf = forge_title.text_frame
    p = tf.paragraphs[0]
    p.text = "MCP CONTEXT FORGE"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Components inside
    components = ["Registry", "Gateway", "Auth & RBAC", "Metrics"]
    for i, comp in enumerate(components):
        add_diagram_box(slide, comp, 1.3 + i * 2, 3.8, 1.7, 0.7, WHITE, IBM_DARK_BLUE)

    add_arrow_down(slide, 4.85, 4.85, 0.4)

    # Bottom layer - Services
    services = ["MCP Servers", "REST APIs", "gRPC", "Legacy"]
    for i, svc in enumerate(services):
        add_diagram_box(slide, svc, 1 + i * 2.2, 5.5, 1.9, 0.7, IBM_LIGHT_GRAY, IBM_DARK_BLUE)

    # ========== SLIDE 13: Enterprise Capabilities ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "ENTERPRISE CAPABILITIES", Inches(0.5))

    caps = [
        ("FEDERATION", ["Aggregate from multiple MCP servers", "Peer gateway discovery", "Redis-backed coordination"]),
        ("MULTI-TENANCY", ["Team-based access control", "Private/Team/Global visibility", "RBAC with roles"]),
        ("SECURITY", ["JWT, OAuth 2.0, SSO", "Multiple IdP support", "30+ security checks"]),
        ("OBSERVABILITY", ["OpenTelemetry integration", "Prometheus metrics", "Distributed tracing"])
    ]

    for i, (title, bullets) in enumerate(caps):
        row = i // 2
        col = i % 2
        x = 0.4 + col * 4.8
        y = 1.5 + row * 2.8

        # Title box
        title_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.5), Inches(0.5))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = IBM_BLUE
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Bullets
        bullets_shape = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.6), Inches(4.3), Inches(2))
        tf = bullets_shape.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(bullets):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(14)
            p.font.color.rgb = IBM_GRAY

    # ========== SLIDE 14: Tool Lifecycle ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "TOOL LIFECYCLE MANAGEMENT", Inches(0.5))

    lifecycle = ["REGISTER", "ENABLE", "DISCOVER", "INVOKE", "AUDIT", "RETIRE"]

    for i, step in enumerate(lifecycle):
        row = i // 3
        col = i % 3
        x = 0.8 + col * 3
        y = 1.8 + row * 2

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = IBM_BLUE if row == 0 else ACCENT_BLUE
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Add arrows between boxes
        if col < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.5), Inches(y + 0.35), Inches(0.4), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = IBM_GRAY

    conclusion = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.6))
    tf = conclusion.text_frame
    p = tf.paragraphs[0]
    p.text = "Full lifecycle visibility and control"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = IBM_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 15: Security Architecture ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "SECURITY ARCHITECTURE", Inches(0.5))

    add_two_column_bullets(
        slide,
        "AUTHENTICATION",
        ["JWT (HS256/RS256)", "OAuth 2.0", "SSO/OIDC", "Basic Auth"],
        "AUTHORIZATION",
        ["RBAC with roles", "Per-resource permissions", "Team isolation"],
        Inches(1.4)
    )

    # IdP section
    idp_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.5))
    tf = idp_title.text_frame
    p = tf.paragraphs[0]
    p.text = "SUPPORTED IDENTITY PROVIDERS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE

    idps = ["GitHub", "Google", "Okta", "Keycloak", "IBM Verify", "MS Entra ID", "Generic OIDC"]
    for i, idp in enumerate(idps):
        x = 0.5 + (i % 4) * 2.3
        y = 4.8 + (i // 4) * 0.7
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2), Inches(0.5))
        box.fill.solid()
        box.fill.fore_color.rgb = IBM_LIGHT_GRAY
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = idp
        p.font.size = Pt(14)
        p.font.color.rgb = IBM_DARK_BLUE
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

    # ========== SLIDE 16: Federation ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "FEDERATION & MULTI-CLUSTER DEPLOYMENT", Inches(0.5))

    # Redis cluster at top
    add_diagram_box(slide, "REDIS CLUSTER\n(Coordination Layer)", 3, 1.5, 4, 0.9, IBM_DARK_BLUE)

    # Three gateways
    gateways = [("Gateway A", "US-East"), ("Gateway B", "EU-West"), ("Gateway C", "AP-South")]
    for i, (name, region) in enumerate(gateways):
        x = 1 + i * 3
        add_diagram_box(slide, f"{name}\n({region})", x, 3, 2.5, 0.9, IBM_BLUE)
        # Connection lines to Redis
        line = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(x + 1.25), Inches(2.4), Inches(0.02), Inches(0.6))
        line.fill.solid()
        line.fill.fore_color.rgb = IBM_GRAY

    # Local MCP servers
    for i in range(3):
        x = 1 + i * 3
        add_diagram_box(slide, "Local MCP\nServers", x, 4.5, 2.5, 0.9, IBM_LIGHT_GRAY, IBM_DARK_BLUE)
        # Connection to gateway
        line = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(x + 1.25), Inches(3.9), Inches(0.02), Inches(0.6))
        line.fill.solid()
        line.fill.fore_color.rgb = IBM_GRAY

    conclusion = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.6))
    tf = conclusion.text_frame
    p = tf.paragraphs[0]
    p.text = "Unified tool discovery across all regions"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 17: Plugin System ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "EXTENSIBLE PLUGIN SYSTEM", Inches(0.5))

    add_two_column_bullets(
        slide,
        "PRE-INVOKE HOOKS",
        ["Input validation", "Access control", "PII detection", "Rate limiting"],
        "POST-INVOKE HOOKS",
        ["Result filtering", "Audit logging", "Response transform", "Metrics capture"],
        Inches(1.4)
    )

    # Built-in plugins
    plugins_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    tf = plugins_title.text_frame
    p = tf.paragraphs[0]
    p.text = "BUILT-IN PLUGINS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE

    plugins = [
        ("PII Filter", "Detects/masks SSN, emails"),
        ("Regex Filter", "Pattern transformations"),
        ("Deny List", "Block specific terms"),
        ("Resource Filter", "Size limits, domain blocking")
    ]

    for i, (name, desc) in enumerate(plugins):
        x = 0.5 + i * 2.3
        box = slide.shapes.add_textbox(Inches(x), Inches(5.1), Inches(2.2), Inches(1))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = IBM_DARK_BLUE
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = IBM_GRAY

    # ========== SLIDE 18: Observability ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "OBSERVABILITY STACK", Inches(0.5))

    # OpenTelemetry at top
    add_diagram_box(slide, "OpenTelemetry", 3.5, 1.5, 3, 0.8, IBM_DARK_BLUE)

    # Three pillars
    pillars = [
        ("TRACES", ["Jaeger", "Zipkin", "Tempo", "Phoenix"]),
        ("METRICS", ["Prometheus", "Grafana"]),
        ("LOGS", ["Structured JSON", "Correlation IDs"])
    ]

    for i, (title, items) in enumerate(pillars):
        x = 0.8 + i * 3.2

        # Title box
        title_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.8), Inches(2.8), Inches(0.6))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = IBM_BLUE
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Items
        items_shape = slide.shapes.add_textbox(Inches(x), Inches(3.5), Inches(2.8), Inches(2))
        tf = items_shape.text_frame
        for j, item in enumerate(items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = IBM_GRAY
            p.alignment = PP_ALIGN.CENTER

    # Additional features
    features = [
        "Correlation IDs for end-to-end tracing",
        "LLM-specific metrics (tokens, costs, latency)",
        "Tool invocation metrics and error rates"
    ]
    add_bullet_points(slide, features, Inches(5.5), font_size=16)

    # ========== SLIDE 19: Deployment ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "FLEXIBLE DEPLOYMENT OPTIONS", Inches(0.5))

    deployments = [
        ("PYPI PACKAGE", ["pip install mcp-contextforge-gateway", "Quick local setup"]),
        ("CONTAINER IMAGES", ["ghcr.io/ibm/mcp-context-forge", "Multi-arch: amd64, arm64, ppc64le"]),
        ("KUBERNETES", ["Helm charts with HPA", "Network policies, RBAC"]),
        ("IBM CLOUD", ["IBM Code Engine", "Serverless, auto-scaling"])
    ]

    for i, (title, items) in enumerate(deployments):
        row = i // 2
        col = i % 2
        x = 0.4 + col * 4.8
        y = 1.5 + row * 2.5

        # Title
        title_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.5), Inches(0.5))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = IBM_BLUE
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Items
        items_shape = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.6), Inches(4.3), Inches(1.5))
        tf = items_shape.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = IBM_GRAY

    # ========== SLIDE 20: Admin UI ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "ADMIN UI: REAL-TIME MANAGEMENT", Inches(0.5))

    # Dashboard mockup
    dashboard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    dashboard.fill.solid()
    dashboard.fill.fore_color.rgb = IBM_LIGHT_GRAY
    dashboard.line.color.rgb = IBM_GRAY

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(0.6))
    header.fill.solid()
    header.fill.fore_color.rgb = IBM_DARK_BLUE

    header_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(5), Inches(0.5))
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = "MCP Context Forge Dashboard"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Stat boxes
    stats = [("Servers", "12"), ("Tools", "47"), ("Resources", "23"), ("Prompts", "8")]
    for i, (label, value) in enumerate(stats):
        x = 0.8 + i * 2.2
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3), Inches(1.9), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE

        value_text = slide.shapes.add_textbox(Inches(x), Inches(2.4), Inches(1.9), Inches(0.6))
        tf = value_text.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = IBM_BLUE
        p.alignment = PP_ALIGN.CENTER

        label_text = slide.shapes.add_textbox(Inches(x), Inches(3), Inches(1.9), Inches(0.4))
        tf = label_text.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = IBM_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Activity section
    activity_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.8), Inches(8.4), Inches(1.8))
    activity_box.fill.solid()
    activity_box.fill.fore_color.rgb = WHITE

    activity_title = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(4), Inches(0.4))
    tf = activity_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Recent Activity"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = IBM_DARK_BLUE

    activities = ["search_database invoked (2s ago)", "new_tool registered (5m ago)", "gateway_b connected (10m ago)"]
    for i, act in enumerate(activities):
        act_text = slide.shapes.add_textbox(Inches(1), Inches(4.4 + i * 0.4), Inches(7.5), Inches(0.4))
        tf = act_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"► {act}"
        p.font.size = Pt(14)
        p.font.color.rgb = IBM_GRAY

    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Built with HTMX + Alpine.js for responsive updates"
    p.font.size = Pt(14)
    p.font.color.rgb = IBM_GRAY
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 21: Demo Introduction ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, IBM_DARK_BLUE)

    demo_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = demo_title.text_frame
    p = tf.paragraphs[0]
    p.text = "LIVE DEMONSTRATION"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    demo_items = [
        "1. MCP Context Forge Admin UI",
        "2. Tool Discovery via API",
        "3. Tool Governance in Action",
        "4. Federation (if time permits)"
    ]

    for i, item in enumerate(demo_items):
        item_text = slide.shapes.add_textbox(Inches(2), Inches(4 + i * 0.6), Inches(6), Inches(0.5))
        tf = item_text.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 22: Key Outcomes ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "KEY OUTCOMES WITH MCP CONTEXT FORGE", Inches(0.5))

    outcomes = [
        ("Custom integrations per tool", "Standardized protocol for all tools"),
        ("Scattered security configurations", "Centralized auth and RBAC"),
        ("No visibility into tool usage", "Complete observability and audit trail"),
        ("Vendor-locked integrations", "Portable across AI providers")
    ]

    # Before/After headers
    before_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4), Inches(0.5))
    tf = before_header.text_frame
    p = tf.paragraphs[0]
    p.text = "BEFORE"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = IBM_GRAY
    p.alignment = PP_ALIGN.CENTER

    after_header = slide.shapes.add_textbox(Inches(5.5), Inches(1.4), Inches(4), Inches(0.5))
    tf = after_header.text_frame
    p = tf.paragraphs[0]
    p.text = "AFTER"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE
    p.alignment = PP_ALIGN.CENTER

    for i, (before, after) in enumerate(outcomes):
        y = 2 + i * 1.1

        # Before box
        before_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(4), Inches(0.9))
        before_box.fill.solid()
        before_box.fill.fore_color.rgb = IBM_LIGHT_GRAY

        before_text = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.2), Inches(3.8), Inches(0.6))
        tf = before_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = before
        p.font.size = Pt(14)
        p.font.color.rgb = IBM_GRAY

        # Arrow
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.6), Inches(y + 0.3), Inches(0.8), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_BLUE

        # After box
        after_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(y), Inches(4), Inches(0.9))
        after_box.fill.solid()
        after_box.fill.fore_color.rgb = ACCENT_BLUE

        after_text = slide.shapes.add_textbox(Inches(5.6), Inches(y + 0.2), Inches(3.8), Inches(0.6))
        tf = after_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = after
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

    # ========== SLIDE 23: Getting Started ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "GETTING STARTED", Inches(0.5))

    # PyPI section
    pypi_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.4))
    tf = pypi_title.text_frame
    p = tf.paragraphs[0]
    p.text = "QUICK START (PyPI)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE

    pypi_code = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(9), Inches(0.8))
    pypi_code.fill.solid()
    pypi_code.fill.fore_color.rgb = IBM_DARK_BLUE

    pypi_text = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(8.6), Inches(0.6))
    tf = pypi_text.text_frame
    p = tf.paragraphs[0]
    p.text = "pip install mcp-contextforge-gateway && mcpgateway serve"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE

    # Docker section
    docker_title = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(0.4))
    tf = docker_title.text_frame
    p = tf.paragraphs[0]
    p.text = "QUICK START (Docker)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE

    docker_code = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    docker_code.fill.solid()
    docker_code.fill.fore_color.rgb = IBM_DARK_BLUE

    docker_text = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(8.6), Inches(0.6))
    tf = docker_text.text_frame
    p = tf.paragraphs[0]
    p.text = "docker run -p 4444:4444 ghcr.io/ibm/mcp-context-forge:latest"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE

    # Resources section
    resources_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.4))
    tf = resources_title.text_frame
    p = tf.paragraphs[0]
    p.text = "RESOURCES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE

    resources = [
        "GitHub: github.com/IBM/mcp-context-forge",
        "Documentation: ibm.github.io/mcp-context-forge",
        "PyPI: pypi.org/project/mcp-contextforge-gateway"
    ]
    add_bullet_points(slide, resources, Inches(5), font_size=18)

    # ========== SLIDE 24: Summary ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    add_title_shape(slide, "SUMMARY", Inches(0.5))

    summary_points = [
        ("1", "MCP provides the STANDARD", "Consistent interface for AI-tool interaction"),
        ("2", "MCP Context Forge provides the CONTROL", "Centralized governance, security, federation"),
        ("3", "Together they enable SCALE", "Production-ready agentic AI infrastructure")
    ]

    for i, (num, title, desc) in enumerate(summary_points):
        y = 1.6 + i * 1.4

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = IBM_BLUE
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Title
        title_text = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(8), Inches(0.5))
        tf = title_text.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = IBM_DARK_BLUE

        # Description
        desc_text = slide.shapes.add_textbox(Inches(1.3), Inches(y + 0.5), Inches(8), Inches(0.5))
        tf = desc_text.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = IBM_GRAY

    # Quote
    quote_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = IBM_LIGHT_GRAY
    quote_box.line.color.rgb = IBM_BLUE

    quote_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.65), Inches(8.6), Inches(0.7))
    tf = quote_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"Tools become governed, reusable platform assets — not one-off integrations."'
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = IBM_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 25: Thank You ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, IBM_DARK_BLUE)

    thank_you = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = thank_you.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.6))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Presenter info
    presenter_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
    tf = presenter_box.text_frame

    p = tf.paragraphs[0]
    p.text = "Ruslan Idelfonso Magana Vsevolodovna"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Data Scientist & AI Engineer"
    p2.font.size = Pt(18)
    p2.font.color.rgb = IBM_LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "IBM Client Innovation Center Italy"
    p3.font.size = Pt(18)
    p3.font.color.rgb = IBM_LIGHT_GRAY
    p3.alignment = PP_ALIGN.CENTER

    # Contact info
    contact = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1))
    tf = contact.text_frame

    p = tf.paragraphs[0]
    p.text = "CONNECT"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "GitHub: github.com/IBM/mcp-context-forge"
    p2.font.size = Pt(16)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Documentation: ibm.github.io/mcp-context-forge"
    p3.font.size = Pt(16)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER

    # Save
    output_path = "/home/user/mcp-context-forge/docs/presentation/MCP-Context-Forge-Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
